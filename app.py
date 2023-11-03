import pptx
import pandas as pd
import shutil, os, sys
from datetime import datetime
import math

# pip3 install pandas python-pptx openpyxl
# shutil should be included with python 3.6 

# # #  setup  # # # # # # # # # #

filename_data = 'input/Youtube_Studio_export.xlsx' #input data
#filename_data = sys.argv[1]
filename_master = "masters/Finii_Reporting_YT.pptx" #path master pptx
output_dir = 'output'

# # #  ende  # # # # # # # # # #

# init
new_report = os.path.splitext(filename_master)[0] + '_edited.pptx'
shutil.copy2(filename_master, new_report)
ppt = pptx.Presentation(new_report)

def insertTextOnSlide(ppt, content, slide_nr, platzhalter_name):

    slide_nr -= 1
    slide = ppt.slides[slide_nr]

    # Den Platzhalter mit dem Namen "xy" suchen
    gesuchter_platzhalter = None

    for platzhalter in slide.placeholders:
        if platzhalter.name == platzhalter_name:
            gesuchter_platzhalter = platzhalter
            break

    # Überprüfen, ob der Platzhalter gefunden wurde
    if gesuchter_platzhalter:
        # Hier kannst du den Platzhalter manipulieren, z.B. Text hinzufügen
        gesuchter_platzhalter.text = content if not content == "nan" else "n/a"
        print(f"Erfolgreich platziert: {platzhalter_name} auf Slide {slide_nr+1} ")

    else:
        print(f"Digger, der Platzhalter mit dem Namen '{platzhalter_name}' auf Slide {slide_nr+1} wurde nicht gefunden.")

    return True

def format_number(number):
    if number >= 1e6:  # Wenn die Zahl eine Million oder mehr ist
        return f'{number / 1e6:.1f}M'  # Formatieren als Millionen (eine Dezimalstelle)
    elif number >= 1e3:  # Wenn die Zahl tausend oder mehr ist
        return f'{number / 1e3:.1f}K'  # Formatieren als Tausende (eine Dezimalstelle)
    else:
        return str(round(number))  # Kleine Zahlen unverändert lassen


def createSlide(ppt, slide_nr, slide_name, tab="Tabellendaten", columns="A:L"):
    
    # load the right part of data
    df = pd.read_excel(filename_data, sheet_name=tab, index_col=None, na_values=['NA'], usecols = columns)
    df["Interactions"] = df["\"Mag ich\"-Bewertungen"] + df["Kommentare hinzugefügt"] + df["Geteilte Inhalte"]
    #print(df)

    # abbruch wenn df leer, zb bei leerem tab 
    # if df.empty:
    #     return True

    df['Veröffentlichungszeitpunkt des Videos'] = pd.to_datetime(df['Veröffentlichungszeitpunkt des Videos'], format="%b %d, %Y")
    df['Date'] = df['Veröffentlichungszeitpunkt des Videos'].dt.strftime("%Y-%m-%d")


    
    print(f"\nDer df (unsortiert) für slide {slide_nr}:\n{df.head()}")

    slide_nr -= 1
    slide = ppt.slides[slide_nr]

    try:

        # top performer ermitteln und in df speichern

        df_sorted = df.sort_values(by='Aufrufe', ascending=False)
        views_top1 = df_sorted.iloc[0] if len(df_sorted) > 0 else None
        views_flop1 = df_sorted.iloc[-1] if len(df_sorted) > 1 else None
        views_average = df_sorted['Aufrufe'].mean()

        df_sorted = df.sort_values(by='Interactions', ascending=False)
        interactions_top1 = df_sorted.iloc[0] if len(df_sorted) > 0 else None
        interactions_flop1 = df_sorted.iloc[-1] if len(df_sorted) > 1 else None
        interactions_average = df_sorted['Interactions'].mean()

        df_sorted = df.sort_values(by='Durchschnittliche Wiedergabedauer in Prozent (%)', ascending=False)
        aufrufe_top1 = df_sorted.iloc[0] if len(df_sorted) > 0 else None
        aufrufe_flop1 = df_sorted.iloc[-1] if len(df_sorted) > 1 else None
        aufrufe_average = df_sorted['Durchschnittliche Wiedergabedauer in Prozent (%)'].mean()


    except:

        pass
        print( f"Ich habe heute leider keinen Top Post für dich auf Slide {slide_nr+1} ({slide_name})")

    # Ab hier beginnt die Zuschreibung zu Slides 

    try:

        if slide_name == "YT_OVER":

            insertTextOnSlide(ppt, f"{format_number(views_average)}", 1, "Inhaltsplatzhalter 21")
            insertTextOnSlide(ppt, f"{format_number(interactions_average)}", 1, "Inhaltsplatzhalter 20")
            insertTextOnSlide(ppt, f"{format_number(aufrufe_average)}%", 1, "Inhaltsplatzhalter 19")

            insertTextOnSlide(ppt, f"{views_top1['Date']}", 1, "Inhaltsplatzhalter 7")
            insertTextOnSlide(ppt, f"{format_number(views_top1['Aufrufe'])}", 1, "Inhaltsplatzhalter 4")
            insertTextOnSlide(ppt, f"{format_number(views_top1['Interactions'])}", 1, "Inhaltsplatzhalter 5")
            insertTextOnSlide(ppt, f"{format_number(views_top1['Durchschnittliche Wiedergabedauer in Prozent (%)'])}%", 1, "Inhaltsplatzhalter 6")
            insertTextOnSlide(ppt, f"{views_top1['Videos']}", 1, "Inhaltsplatzhalter 8")

            insertTextOnSlide(ppt, f"{interactions_top1['Date']}", 1, "Inhaltsplatzhalter 12")
            insertTextOnSlide(ppt, f"{format_number(interactions_top1['Aufrufe'])}", 1, "Inhaltsplatzhalter 9")
            insertTextOnSlide(ppt, f"{format_number(interactions_top1['Interactions'])}", 1, "Inhaltsplatzhalter 10")
            insertTextOnSlide(ppt, f"{format_number(interactions_top1['Durchschnittliche Wiedergabedauer in Prozent (%)'])}%", 1, "Inhaltsplatzhalter 11")
            insertTextOnSlide(ppt, f"{interactions_top1['Videos']}", 1, "Inhaltsplatzhalter 13")

            insertTextOnSlide(ppt, f"{aufrufe_top1['Date']}", 1, "Inhaltsplatzhalter 17")
            insertTextOnSlide(ppt, f"{format_number(aufrufe_top1['Aufrufe'])}", 1, "Inhaltsplatzhalter 14")
            insertTextOnSlide(ppt, f"{format_number(aufrufe_top1['Interactions'])}", 1, "Inhaltsplatzhalter 15")
            insertTextOnSlide(ppt, f"{format_number(aufrufe_top1['Durchschnittliche Wiedergabedauer in Prozent (%)'])}%", 1, "Inhaltsplatzhalter 16")
            insertTextOnSlide(ppt, f"{aufrufe_top1['Videos']}", 1, "Inhaltsplatzhalter 18")


        elif slide_name == "YT_UNDER":

            insertTextOnSlide(ppt, f"{format_number(views_average)}", 2, "Inhaltsplatzhalter 21")
            insertTextOnSlide(ppt, f"{format_number(interactions_average)}", 2, "Inhaltsplatzhalter 20")
            insertTextOnSlide(ppt, f"{format_number(aufrufe_average)}%", 2, "Inhaltsplatzhalter 19")

            insertTextOnSlide(ppt, f"{views_flop1['Date']}", 2, "Inhaltsplatzhalter 7")
            insertTextOnSlide(ppt, f"{format_number(views_flop1['Aufrufe'])}", 2, "Inhaltsplatzhalter 4")
            insertTextOnSlide(ppt, f"{format_number(views_flop1['Interactions'])}", 2, "Inhaltsplatzhalter 5")
            insertTextOnSlide(ppt, f"{format_number(views_flop1['Durchschnittliche Wiedergabedauer in Prozent (%)'])}%", 2, "Inhaltsplatzhalter 6")
            insertTextOnSlide(ppt, f"{views_flop1['Videos']}", 2, "Inhaltsplatzhalter 8")

            insertTextOnSlide(ppt, f"{interactions_flop1['Date']}", 2, "Inhaltsplatzhalter 12")
            insertTextOnSlide(ppt, f"{format_number(interactions_flop1['Aufrufe'])}", 2, "Inhaltsplatzhalter 9")
            insertTextOnSlide(ppt, f"{format_number(interactions_flop1['Interactions'])}", 2, "Inhaltsplatzhalter 10")
            insertTextOnSlide(ppt, f"{format_number(interactions_flop1['Durchschnittliche Wiedergabedauer in Prozent (%)'])}%", 2, "Inhaltsplatzhalter 11")
            insertTextOnSlide(ppt, f"{interactions_flop1['Videos']}", 2, "Inhaltsplatzhalter 13")

            insertTextOnSlide(ppt, f"{aufrufe_flop1['Date']}", 2, "Inhaltsplatzhalter 17")
            insertTextOnSlide(ppt, f"{format_number(aufrufe_flop1['Aufrufe'])}", 2, "Inhaltsplatzhalter 14")
            insertTextOnSlide(ppt, f"{format_number(aufrufe_flop1['Interactions'])}", 2, "Inhaltsplatzhalter 15")
            insertTextOnSlide(ppt, f"{format_number(aufrufe_flop1['Durchschnittliche Wiedergabedauer in Prozent (%)'])}%", 2, "Inhaltsplatzhalter 16")
            insertTextOnSlide(ppt, f"{aufrufe_flop1['Videos']}", 2, "Inhaltsplatzhalter 18")


        else:

            pass

    except:
        
        pass
        print(f"! ! ! ! Abgebrochen: Diese Folie ({slide_name} auf Slide {slide_nr+1}) nicht komplett eingesetzt ! ! ! !")
        

# Control Panel
createSlide(ppt, 1, "YT_OVER") # Presentation object, platform, slide nr, slide_name
createSlide(ppt, 2, "YT_UNDER") # Presentation object, platform, slide nr, slide_name

# saving file
ppt.save(new_report)

# move the copied file to the output directory
new_file_path = os.path.join(output_dir, os.path.basename(new_report))
shutil.move(new_report, new_file_path)
